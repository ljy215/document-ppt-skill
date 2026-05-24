#!/usr/bin/env python3
"""Compile document-ppt Slide-JSON into a native PPTX deck."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


SLIDE_WIDE = (13.333333, 7.5)
SLIDE_4_3 = (10.0, 7.5)


def load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def hex_color(value: str, fallback: str) -> RGBColor:
    raw = (value or fallback).lstrip("#")
    if len(raw) != 6:
        raw = fallback.lstrip("#")
    return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def validate_deck(deck: dict[str, Any]) -> None:
    if deck.get("schema_version") != "document-ppt.slide.v1":
        raise ValueError("schema_version must be document-ppt.slide.v1")
    if not isinstance(deck.get("slides"), list) or not deck["slides"]:
        raise ValueError("slides must be a non-empty array")


def resolve_asset(deck_path: Path, deck: dict[str, Any], asset_url: str) -> Path:
    raw = Path(asset_url)
    if raw.is_absolute():
        return raw
    source_manifest = deck.get("deck", {}).get("source_manifest")
    if source_manifest:
        manifest_dir = Path(source_manifest).expanduser()
        if not manifest_dir.is_absolute():
            manifest_dir = (deck_path.parent / manifest_dir).resolve()
        else:
            manifest_dir = manifest_dir.resolve()
        if manifest_dir.name == "manifest.json":
            manifest_dir = manifest_dir.parent
        candidate = manifest_dir / raw
        if candidate.exists():
            return candidate
    return (deck_path.parent / raw).resolve()


def set_page_size(prs: Presentation, aspect_ratio: str) -> tuple[float, float]:
    width, height = SLIDE_4_3 if aspect_ratio == "4:3" else SLIDE_WIDE
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)
    return width, height


def add_textbox(slide: Any, text: str, x: float, y: float, w: float, h: float, font_size: int, color: RGBColor, bold: bool = False) -> Any:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_bullets(slide: Any, bullets: list[dict[str, Any]], x: float, y: float, w: float, h: float, color: RGBColor, accent: RGBColor) -> None:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet.get("text", "")
        paragraph.level = 0
        paragraph.space_after = Pt(8)
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        run.font.size = Pt(20 if len(bullets) <= 4 else 17)
        run.font.color.rgb = accent if bullet.get("emphasis") in ("strong", "evidence") else color
        run.font.bold = bullet.get("emphasis") == "strong"


def add_visual(slide: Any, deck_path: Path, deck: dict[str, Any], visual: dict[str, Any], slide_w: float, slide_h: float, color: RGBColor) -> None:
    layout = visual.get("layout") or {"x": 0.55, "y": 0.22, "w": 0.36, "h": 0.56}
    image_path = resolve_asset(deck_path, deck, visual.get("asset_url", ""))
    x = slide_w * float(layout.get("x", 0))
    y = slide_h * float(layout.get("y", 0))
    w = slide_w * float(layout.get("w", 0.4))
    h = slide_h * float(layout.get("h", 0.4))
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        add_textbox(slide, f"Missing asset:\n{visual.get('asset_url', '')}", x, y, w, h, 12, color)
    caption = visual.get("caption")
    if caption:
        add_textbox(slide, caption[:160], x, min(slide_h - 0.45, y + h + 0.06), w, 0.35, 9, color)


def set_background(slide: Any, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_transition(slide: Any, transition: str | None) -> None:
    if not transition or transition == "none":
        return
    transition_map = {"fade": "fade", "wipe": "wipe", "push": "push"}
    tag = transition_map.get(transition)
    if not tag:
        return
    slide_el = slide._element
    existing = slide_el.find("{http://schemas.openxmlformats.org/presentationml/2006/main}transition")
    if existing is not None:
        slide_el.remove(existing)
    transition_el = OxmlElement("p:transition")
    transition_el.set("spd", "med")
    transition_el.append(OxmlElement(f"p:{tag}"))
    slide_el.insert(0, transition_el)


def add_notes(slide: Any, slide_data: dict[str, Any]) -> None:
    notes_lines = []
    if slide_data.get("speaker_notes"):
        notes_lines.append(slide_data["speaker_notes"])
    animation_lines = []
    for bullet in slide_data.get("bullets") or []:
        anim = bullet.get("animation") or {}
        animation_lines.append(f"Bullet animation: {anim.get('type', 'none')} order={anim.get('order', 0)} text={bullet.get('text', '')}")
    for visual in slide_data.get("visuals") or []:
        anim = visual.get("animation") or {}
        animation_lines.append(f"Visual animation: {anim.get('type', 'none')} order={anim.get('order', 0)} asset={visual.get('asset_id', '')}")
    if animation_lines:
        notes_lines.append("Animation intent:\n" + "\n".join(animation_lines))
    if not notes_lines:
        return
    try:
        notes_frame = slide.notes_slide.notes_text_frame
        notes_frame.text = "\n\n".join(notes_lines)
    except Exception:
        pass


def layout_regions(layout: str, has_visual: bool, slide_w: float, slide_h: float) -> tuple[tuple[float, float, float, float], tuple[float, float, float, float] | None]:
    if layout == "visual-left" and has_visual:
        return (7.15, 1.72, 5.25, 4.35), (0.65, 1.7, 5.75, 4.35)
    if layout == "visual-full" and has_visual:
        return (0.85, 6.05, 11.6, 0.9), (0.85, 1.55, 11.6, 4.3)
    if layout == "title":
        return (1.0, 4.7, 11.2, 1.1), None
    if has_visual:
        return (0.85, 1.75, 5.75, 4.65), (7.05, 1.65, 5.35, 4.45)
    return (1.0, 1.9, 11.3, 4.8), None


def render_slide(prs: Presentation, deck_path: Path, deck: dict[str, Any], slide_data: dict[str, Any], slide_w: float, slide_h: float) -> None:
    theme = deck.get("deck", {}).get("theme", {})
    palette = theme.get("palette", {})
    bg = hex_color(palette.get("background"), "#F8FAFC")
    fg = hex_color(palette.get("foreground"), "#111827")
    accent = hex_color(palette.get("accent"), "#2563EB")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, bg)
    add_transition(slide, slide_data.get("transition"))

    layout = slide_data.get("layout", "bullets")
    title_size = 34 if layout != "title" else 42
    title_y = 0.45 if layout != "title" else 1.85
    add_textbox(slide, slide_data.get("title", ""), 0.78, title_y, slide_w - 1.55, 0.65, title_size, fg, bold=True)
    if slide_data.get("subtitle"):
        add_textbox(slide, slide_data["subtitle"], 0.82, title_y + 0.72, slide_w - 1.65, 0.45, 17, accent)

    visuals = slide_data.get("visuals") or []
    bullet_region, visual_region = layout_regions(layout, bool(visuals), slide_w, slide_h)
    if slide_data.get("bullets"):
        add_bullets(slide, slide_data["bullets"], *bullet_region, color=fg, accent=accent)

    for index, visual in enumerate(visuals[:3]):
        if visual_region:
            x, y, w, h = visual_region
            if len(visuals) > 1:
                h = (visual_region[3] - 0.18 * (len(visuals) - 1)) / len(visuals)
                y = visual_region[1] + index * (h + 0.18)
            visual = dict(visual)
            visual["layout"] = {"x": x / slide_w, "y": y / slide_h, "w": w / slide_w, "h": h / slide_h}
        add_visual(slide, deck_path, deck, visual, slide_w, slide_h, fg)

    add_notes(slide, slide_data)


def compile_pptx(deck_path: Path, output_path: Path) -> dict[str, Any]:
    deck = load_json(deck_path)
    validate_deck(deck)
    prs = Presentation()
    aspect = deck.get("deck", {}).get("theme", {}).get("aspect_ratio", "16:9")
    slide_w, slide_h = set_page_size(prs, aspect)
    for slide_data in deck["slides"]:
        render_slide(prs, deck_path, deck, slide_data, slide_w, slide_h)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return {"output": str(output_path), "slides": len(deck["slides"])}


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("slide_json", help="Path to document-ppt.slide.v1 JSON")
    parser.add_argument("--output", required=True, help="Output PPTX path")
    args = parser.parse_args()
    result = compile_pptx(Path(args.slide_json).resolve(), Path(args.output).resolve())
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
